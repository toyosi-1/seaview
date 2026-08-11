from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io

# ── Brand colours ──────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x8A)   # #1e3a8a
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)   # #dbeafe
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLATE       = RGBColor(0x47, 0x56, 0x69)   # #475569
DARK        = RGBColor(0x0F, 0x17, 0x2A)   # near-black
GREEN       = RGBColor(0x16, 0xA3, 0x4A)   # #16a34a
AMBER       = RGBColor(0xD9, 0x77, 0x06)   # #d97706

W = Inches(13.33)   # 16:9 widescreen width
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # fully blank layout

# ── Helper utilities ────────────────────────────────────────────
def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def label(slide, text, left, top, width, height,
          font_size=Pt(12), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def pill(slide, text, left, top, fill, font_color=WHITE, font_size=Pt(9.5)):
    w = Inches(1.6)
    h = Inches(0.32)
    b = box(slide, left, top, w, h, fill_color=fill)
    b.line.fill.background()
    tf = b.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = True
    run.font.color.rgb = font_color

def section_header(slide, title, subtitle=""):
    """Dark navy full-width header bar"""
    box(slide, 0, 0, W, Inches(1.4), fill_color=NAVY)
    label(slide, title,
          Inches(0.5), Inches(0.18), Inches(10), Inches(0.7),
          font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        label(slide, subtitle,
              Inches(0.5), Inches(0.82), Inches(10), Inches(0.45),
              font_size=Pt(13), color=RGBColor(0xBA, 0xD5, 0xFD))

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, DARK)

# Full-bleed gradient effect — two overlapping boxes
box(s1, 0, 0, W, H, fill_color=NAVY)
box(s1, 0, Inches(4.5), W, Inches(3.0), fill_color=DARK)

# Diagonal accent stripe
acc = s1.shapes.add_shape(1, Inches(9.8), 0, Inches(3.53), H)
acc.fill.solid()
acc.fill.fore_color.rgb = RGBColor(0x1D, 0x4E, 0xD8)
acc.line.fill.background()

# Logo placeholder circle
circ = s1.shapes.add_shape(9,  # oval
    Inches(10.3), Inches(1.6), Inches(2.4), Inches(2.4))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x2A, 0x5F, 0xCC)
circ.line.color.rgb = WHITE
circ.line.width = Pt(2)
tf = circ.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "SVP"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = WHITE

# Tag
label(s1, "PROPOSAL FOR FEDERAL GOVERNMENT DIGITALIZATION MANDATE",
      Inches(0.5), Inches(1.3), Inches(9), Inches(0.4),
      font_size=Pt(10), bold=True,
      color=RGBColor(0x93, 0xC5, 0xFD))

# Main title
label(s1, "Digital Procurement &\nContractor Management\nPortal",
      Inches(0.5), Inches(1.75), Inches(8.8), Inches(2.8),
      font_size=Pt(40), bold=True, color=WHITE)

# Subtitle
label(s1, "In alignment with the Federal Government's directive to digitalize all parastatal operations",
      Inches(0.5), Inches(4.55), Inches(9.2), Inches(0.6),
      font_size=Pt(14), color=RGBColor(0xBA, 0xD5, 0xFD))

# Company
label(s1, "Sea View Properties Ltd  \u00b7  A Subsidiary of the Nigerian Ports Authority",
      Inches(0.5), Inches(5.3), Inches(8.8), Inches(0.5),
      font_size=Pt(11), color=RGBColor(0x7D, 0xA8, 0xE8))

# Date
label(s1, "July 2026",
      Inches(0.5), Inches(6.6), Inches(3), Inches(0.4),
      font_size=Pt(11), color=RGBColor(0x93, 0xC5, 0xFD))


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — POLICY & REGULATORY CONTEXT
# ══════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, RGBColor(0xF8, 0xFA, 0xFF))
section_header(s2, "Policy & Regulatory Context", "The national framework driving digitalization of government operations")

# NITDA quote box
quote_box = box(s2, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5),
               fill_color=NAVY)
quote_box.line.fill.background()
label(s2, "\u201cI am aware that [agencies] are working on how to use emerging technologies to come up with new business processes that will improve service delivery.\u201d",
      Inches(0.8), Inches(1.75), Inches(11.5), Inches(0.8),
      font_size=Pt(13), color=WHITE)
label(s2, "\u2014 Malam Kashifu Inuwa Abdullahi, CCIE  \u00b7  Director General, NITDA",
      Inches(0.8), Inches(2.55), Inches(11.5), Inches(0.4),
      font_size=Pt(11), bold=True, color=RGBColor(0xBA, 0xD5, 0xFD))

# Policy framework cards
policies = [
    ("NDEPS", "National Digital Economy Policy & Strategy",
     "Federal framework mandating all MDAs & parastatals to adopt digital technologies"),
    ("NITDA", "Digital Transformation Directive",
     "DG NITDA advocates paperless operations for \u201coperational excellence, cost saving & efficiency\u201d"),
    ("BPP", "E-Procurement Mandate",
     "Bureau of Public Procurement pushes e-procurement under the Public Procurement Act 2007"),
    ("Tinubu", "Digital Economy Agenda",
     "Administration prioritizes transparency, anti-corruption & reduced cost of governance"),
]

pol_w = Inches(2.9)
pol_h = Inches(2.6)
pol_gap = Inches(0.15)
pol_y = Inches(3.4)
for i, (tag, title, desc) in enumerate(policies):
    px = Inches(0.5) + i * (pol_w + pol_gap)
    box(s2, px, pol_y, pol_w, pol_h, fill_color=WHITE,
        line_color=RGBColor(0xBF, 0xDB, 0xFE), line_width=Pt(1))
    box(s2, px, pol_y, pol_w, Inches(0.5), fill_color=NAVY)
    label(s2, tag,
          px + Inches(0.1), pol_y + Inches(0.05), pol_w - Inches(0.2), Inches(0.4),
          font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s2, title,
          px + Inches(0.15), pol_y + Inches(0.65), pol_w - Inches(0.3), Inches(0.7),
          font_size=Pt(11), bold=True, color=NAVY)
    label(s2, desc,
          px + Inches(0.15), pol_y + Inches(1.35), pol_w - Inches(0.3), Inches(1.1),
          font_size=Pt(9.5), color=SLATE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, WHITE)
section_header(s3, "The Problem", "Current challenges & why digitalization is a strategic imperative")

problems = [
    ("📋", "Paper-Based Process",
     "Proposals, contracts and approvals are managed on paper — slow, error-prone, and hard to track."),
    ("🔍", "No Audit Trail",
     "There is no centralised record of who approved what, when, and why — creating compliance risk."),
    ("⏳", "Approval Bottlenecks",
     "Multi-stage approvals happen via email and phone calls, causing delays and lost correspondence."),
    ("💳", "Payment Delays",
     "Contractor payments are processed manually with no structured verification or payment tracking."),
    ("🔒", "No Access Control",
     "Sensitive financial and procurement data is not protected by role-based permissions."),
]

cols = 2
card_w = Inches(5.8)
card_h = Inches(1.4)
x_start = Inches(0.4)
y_start = Inches(1.6)
gap_x = Inches(0.5)
gap_y = Inches(0.22)

for i, (icon, title, desc) in enumerate(problems):
    row = i // cols
    col = i % cols
    if i == 4:   # last card centred
        x = Inches(3.77)
    else:
        x = x_start + col * (card_w + gap_x)
    y = y_start + row * (card_h + gap_y)

    card = box(s3, x, y, card_w, card_h,
               fill_color=LIGHT_BLUE,
               line_color=RGBColor(0xBF, 0xDB, 0xFE), line_width=Pt(1))

    # Left accent
    box(s3, x, y, Inches(0.07), card_h, fill_color=NAVY)

    label(s3, f"{icon}  {title}",
          x + Inches(0.18), y + Inches(0.12), card_w - Inches(0.3), Inches(0.4),
          font_size=Pt(12), bold=True, color=NAVY)
    label(s3, desc,
          x + Inches(0.18), y + Inches(0.52), card_w - Inches(0.3), Inches(0.8),
          font_size=Pt(10.5), color=SLATE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — THE SOLUTION  (architecture overview)
# ══════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, RGBColor(0xF8, 0xFA, 0xFF))
section_header(s4, "The Solution", "A centralised, role-based digital procurement platform")

# Central hub
hub_x, hub_y = Inches(5.3), Inches(2.9)
hub_r = Inches(1.3)
hub = s4.shapes.add_shape(9, hub_x, hub_y, hub_r * 2, hub_r)
hub.fill.solid()
hub.fill.fore_color.rgb = NAVY
hub.line.fill.background()
tf = hub.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "Sea View\nPortal"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = WHITE

modules = [
    (Inches(0.5),  Inches(2.4),  "Contractor\nOnboarding",   RGBColor(0x2A,0x77,0xDD)),
    (Inches(0.5),  Inches(4.3),  "Proposal\nSubmission",     RGBColor(0x16,0x96,0x8A)),
    (Inches(3.2),  Inches(1.5),  "Approval\nWorkflow",       RGBColor(0x7C,0x3A,0xED)),
    (Inches(3.2),  Inches(5.2),  "Contract\nAward & PDF",    RGBColor(0xD9,0x77,0x06)),
    (Inches(9.8),  Inches(2.4),  "Audit &\nCompliance",      RGBColor(0xDC,0x26,0x26)),
    (Inches(9.8),  Inches(4.3),  "Payment\nProcessing",      RGBColor(0x05,0x96,0x69)),
    (Inches(6.5),  Inches(1.5),  "Dashboard &\nReports",     RGBColor(0x0E,0x7A,0xCA)),
    (Inches(6.5),  Inches(5.2),  "Notifications\n& Alerts",  RGBColor(0xBE,0x18,0x5D)),
]

for mx, my, mtitle, mcol in modules:
    mw, mh = Inches(2.4), Inches(0.9)
    mb = box(s4, mx, my, mw, mh, fill_color=mcol)
    mb.line.fill.background()
    tf = mb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = mtitle
    r.font.size  = Pt(11)
    r.font.bold  = True
    r.font.color.rgb = WHITE


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY WORKFLOW
# ══════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5, WHITE)
section_header(s5, "End-to-End Procurement Workflow", "From contractor registration to payment — fully automated")

steps = [
    ("1", "Contractor\nRegisters",     RGBColor(0x2A,0x77,0xDD), "CAC, TIN, banking docs uploaded & verified"),
    ("2", "Submits\nProposal",         RGBColor(0x7C,0x3A,0xED), "Project scope, cost estimate & documents"),
    ("3", "MD Initial\nReview",        RGBColor(0x0E,0x7A,0xCA), "MD approves or returns for clarification"),
    ("4", "Procurement\nAppraisal",    RGBColor(0x05,0x96,0x69), "Officer appraises, Head of Procurement signs off"),
    ("5", "MD Final\nApproval",        NAVY,                      "Managing Director gives final go-ahead"),
    ("6", "Contract\nAwarded",         RGBColor(0xD9,0x77,0x06), "PDF award letter generated & signed digitally"),
    ("7", "Completion\nReport",        RGBColor(0x16,0x96,0x8A), "Contractor submits with images & certificates"),
    ("8", "Audit &\nPayment",          RGBColor(0xDC,0x26,0x26), "Audit review → Accounts processes payment"),
]

step_w = Inches(1.42)
step_h = Inches(2.3)
gap    = Inches(0.12)
y_top  = Inches(1.55)
x0     = Inches(0.28)

for i, (num, title, col, note) in enumerate(steps):
    x = x0 + i * (step_w + gap)

    # Arrow connector (except last)
    if i < len(steps) - 1:
        arr = box(s5, x + step_w, y_top + Inches(0.75),
                  gap + Inches(0.02), Inches(0.1),
                  fill_color=RGBColor(0xCB, 0xD5, 0xE1))

    # Card
    cb = box(s5, x, y_top, step_w, step_h, fill_color=col)
    cb.line.fill.background()

    # Number badge
    nb = s5.shapes.add_shape(9, x + Inches(0.47), y_top + Inches(0.15),
                              Inches(0.5), Inches(0.5))
    nb.fill.solid()
    nb.fill.fore_color.rgb = WHITE
    nb.line.fill.background()
    tf = nb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r  = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = col

    # Title
    label(s5, title,
          x + Inches(0.08), y_top + Inches(0.72),
          step_w - Inches(0.16), Inches(0.85),
          font_size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Note
    label(s5, note,
          x + Inches(0.08), y_top + Inches(1.56),
          step_w - Inches(0.16), Inches(0.65),
          font_size=Pt(8.5), color=RGBColor(0xE2,0xE8,0xF0),
          align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — PLATFORM SCREENSHOTS
# ══════════════════════════════════════════════════════════════════
SCR = "/Users/kingtoy/Documents/Procurement & Contractor Management Portal/CascadeProjects/windsurf-project/screenshots"

screen_slides = [
    ("01_dashboard.png",         "Executive Dashboard",       "Real-time KPIs, recent proposals, and activity feed — tailored per user role"),
    ("02_contractors.png",       "Contractor Registry",       "All registered contractors with status badges, contact details and quick access"),
    ("03_contractor_profile.png","Contractor Profile",        "Full company info, banking details, documents, and staff action controls"),
    ("04_proposal_workflow.png", "Proposal & Approval Panel", "Proposal details, workflow timeline, and role-based action buttons for approvers"),
    ("05_payments.png",          "Payment Management",        "Payment pipeline with amounts, status tracking, and one-click review access"),
]

for img_file, title, subtitle in screen_slides:
    ss = prs.slides.add_slide(BLANK)
    bg(ss, RGBColor(0x0F, 0x17, 0x2A))

    # Top label bar
    box(ss, 0, 0, W, Inches(0.9), fill_color=NAVY)
    label(ss, title,
          Inches(0.5), Inches(0.1), Inches(9), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)
    label(ss, subtitle,
          Inches(0.5), Inches(0.54), Inches(9), Inches(0.3),
          font_size=Pt(11), color=RGBColor(0xBA, 0xD5, 0xFD))

    # Screenshot image with shadow effect
    shadow = box(ss, Inches(0.55), Inches(1.08), Inches(12.25), Inches(6.08),
                 fill_color=RGBColor(0x05, 0x0A, 0x1A))
    pic_path = f"{SCR}/{img_file}"
    ss.shapes.add_picture(pic_path, Inches(0.45), Inches(0.98), Inches(12.25), Inches(6.08))

    # Slide number tag
    label(ss, f"Sea View Properties  ·  Procurement Portal",
          Inches(0.5), Inches(7.22), Inches(8), Inches(0.25),
          font_size=Pt(8), color=RGBColor(0x47, 0x56, 0x69))


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — USER ROLES & ACCESS
# ══════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6, RGBColor(0xF8, 0xFA, 0xFF))
section_header(s6, "User Roles & Access Control", "Every user sees only what they need — enforced at the database level")

roles = [
    ("Contractor",           "Register · Submit proposals · Upload completions · Track payments",       RGBColor(0x2A,0x77,0xDD)),
    ("Managing Director",    "Initial & final proposal approval · Contractor verification · Dashboard", NAVY),
    ("Procurement Officer",  "Appraise proposals · Forward to Head of Procurement",                     RGBColor(0x7C,0x3A,0xED)),
    ("Head of Procurement",  "Final procurement review · Reject or approve for MD sign-off",            RGBColor(0x05,0x96,0x69)),
    ("Head of Audit",        "Review completion reports · Forward to accounts after verification",      RGBColor(0xDC,0x26,0x26)),
    ("Head of Accounts",     "Process payments · Upload bank transfer evidence · Manage deductions",    RGBColor(0xD9,0x77,0x06)),
    ("ICT Administrator",    "User management · System audit logs · Role assignments",                  SLATE),
]

rw  = Inches(5.5)
rh  = Inches(0.66)
gy  = Inches(0.12)
rx0 = Inches(0.45)
ry0 = Inches(1.55)

for i, (role, access, col) in enumerate(roles):
    col_idx = i % 2
    row_idx = i // 2
    if i == 6:  # last centred
        rx = Inches(3.95)
        ry = ry0 + 3 * (rh + gy)
    else:
        rx = rx0 + col_idx * (rw + Inches(0.43))
        ry = ry0 + row_idx * (rh + gy)

    box(s6, rx, ry, rw, rh, fill_color=WHITE,
        line_color=RGBColor(0xE2,0xE8,0xF0), line_width=Pt(1))
    # Left coloured bar
    box(s6, rx, ry, Inches(0.07), rh, fill_color=col)

    label(s6, role,
          rx + Inches(0.18), ry + Inches(0.06), rw - Inches(0.25), Inches(0.3),
          font_size=Pt(11), bold=True, color=NAVY)
    label(s6, access,
          rx + Inches(0.18), ry + Inches(0.35), rw - Inches(0.25), Inches(0.28),
          font_size=Pt(9.5), color=SLATE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7, WHITE)
section_header(s7, "Key Features", "Built for compliance, transparency and speed")

features = [
    ("🔐", "Role-Based Access Control",  "Supabase Row Level Security — data is filtered at the database, not just the UI."),
    ("📄", "PDF Award Letters",           "Auto-generated, digitally signed award letters with MD's uploaded signature."),
    ("📊", "Executive Dashboard",         "Real-time KPIs: proposals pending, contracts active, payments due — all in one view."),
    ("🔔", "Real-Time Notifications",     "Instant alerts for every workflow event — approvals, rejections, awards, payments."),
    ("🛡️", "Immutable Audit Log",         "Every action is time-stamped and actor-attributed. Records cannot be edited or deleted."),
    ("🔎", "Universal Search",            "Find any contractor, proposal, contract or payment instantly from the top bar."),
    ("📱", "Mobile Responsive",           "Works on phones, tablets and desktops — accessible from anywhere."),
    ("�", "Seamless Integrations",         "Built on open standards — integrates with existing systems, email and document workflows."),
]

fw = Inches(5.7)
fh = Inches(1.18)
fgx = Inches(0.5)
fgy = Inches(0.16)
fx0 = Inches(0.4)
fy0 = Inches(1.55)

for i, (icon, feat, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    fx = fx0 + col * (fw + fgx)
    fy = fy0 + row * (fh + fgy)

    box(s7, fx, fy, fw, fh, fill_color=LIGHT_BLUE,
        line_color=RGBColor(0xBF,0xDB,0xFE), line_width=Pt(1))
    box(s7, fx, fy, Inches(0.07), fh, fill_color=NAVY)

    label(s7, f"{icon}  {feat}",
          fx + Inches(0.2), fy + Inches(0.1), fw - Inches(0.28), Inches(0.38),
          font_size=Pt(12), bold=True, color=NAVY)
    label(s7, desc,
          fx + Inches(0.2), fy + Inches(0.5), fw - Inches(0.28), Inches(0.62),
          font_size=Pt(10.5), color=SLATE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — CALL TO ACTION
# ══════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8, DARK)
box(s8, 0, 0, W, H, fill_color=NAVY)

# Accent panel
box(s8, Inches(9.5), 0, Inches(3.83), H, fill_color=RGBColor(0x1D,0x4E,0xD8))

label(s8, "Ready to Launch",
      Inches(0.6), Inches(1.2), Inches(8.5), Inches(1.0),
      font_size=Pt(42), bold=True, color=WHITE)

label(s8, "The platform is built and production-ready.\nWe need approval to connect it to our Supabase account\nand go live within 48 hours.",
      Inches(0.6), Inches(2.4), Inches(8.5), Inches(1.5),
      font_size=Pt(15), color=RGBColor(0xBA,0xD5,0xFD))

next_steps = [
    ("48 hrs",  "Connect Supabase credentials & go live"),
    ("Week 1",  "Onboard ICT Admin + create staff accounts"),
    ("Week 2",  "Contractor registration & first proposals"),
    ("Month 1", "Full procurement cycle running on the platform"),
]

for i, (when, what) in enumerate(next_steps):
    ny = Inches(4.1) + i * Inches(0.65)
    box(s8, Inches(0.6), ny, Inches(1.1), Inches(0.48),
        fill_color=GREEN)
    label(s8, when,
          Inches(0.6), ny + Inches(0.05), Inches(1.1), Inches(0.38),
          font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s8, what,
          Inches(1.85), ny + Inches(0.06), Inches(7.0), Inches(0.38),
          font_size=Pt(12), color=WHITE)

label(s8, "Sea View Properties Ltd  \u00b7  Confidential  \u00b7  July 2026",
      Inches(0.6), Inches(6.9), Inches(8.5), Inches(0.4),
      font_size=Pt(10), color=RGBColor(0x7D,0xA8,0xE8))

# Right panel content
label(s8, "Contact",
      Inches(9.8), Inches(1.8), Inches(3.0), Inches(0.45),
      font_size=Pt(14), bold=True, color=WHITE)

contact_lines = [
    "ICT Department",
    "Sea View Properties Ltd",
    "",
    "📍 Nigerian Ports Authority",
    "    Subsidiary",
    "",
    "🌐 Platform: Ready",
    "⏱  Go-Live: 48 hours",
]
for i, line in enumerate(contact_lines):
    label(s8, line,
          Inches(9.8), Inches(2.35) + i * Inches(0.42),
          Inches(3.1), Inches(0.4),
          font_size=Pt(10.5), color=RGBColor(0xBA,0xD5,0xFD))


# ── Save ─────────────────────────────────────────────────────────
out = "/Users/kingtoy/Documents/Procurement & Contractor Management Portal/CascadeProjects/windsurf-project/SeaView_Procurement_Portal_Proposal.pptx"
prs.save(out)
print(f"✅  Saved → {out}")
